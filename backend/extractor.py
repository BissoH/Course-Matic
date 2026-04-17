# Dispatches to the correct parser based on file extension and returns the extracted plain text.
# Each branch uses a pure-Python library to avoid system-level dependencies that would complicate the Docker build.

import os
from pypdf import PdfReader
import docx
from pptx import Presentation

def extract_text(file_path: str) -> str:
    """
    Takes a file path, detects the extension, and returns the extracted text.
    Returns an empty string on missing files, unsupported extensions, or extraction errors so the caller can handle the empty case uniformly.
    """
    if not os.path.exists(file_path):
        return ""

    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    text = ""

    try:
        # pypdf is used for PDF extraction because it is pure Python with no system dependencies.
        if file_extension == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                extracted = page.extract_text()
                # Pages that consist only of images return None, so the guard prevents a TypeError during concatenation.
                if extracted:
                    text += extracted + "\n"

        # python-docx reads paragraph text only; inline tables and headers are ignored, which is acceptable for typical lecture notes.
        elif file_extension == '.docx':
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"

        # python-pptx exposes per-slide shapes; the hasattr check skips media elements such as images and charts that contain no textual content.
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        # Plain text files are read directly with UTF-8 decoding to handle non-ASCII characters commonly found in academic material.
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

        else:
            print(f"Unsupported file type: {file_extension}")
            return ""

        # Trailing whitespace is stripped so downstream word counts and truncation operate on clean text.
        return text.strip()

    except Exception as e:
        # Any parser error results in an empty string rather than a raised exception, allowing the upload endpoint to remain responsive even when a file is corrupted.
        print(f"Error extracting text from {file_path}: {e}")
        return ""
